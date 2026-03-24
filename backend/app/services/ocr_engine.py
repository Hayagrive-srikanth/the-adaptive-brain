import os
import logging
import tempfile
from typing import Optional
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
from app.config import settings

logger = logging.getLogger(__name__)


def process_document(file_path: str, file_type: str) -> str:
    """Process a document and extract text content."""
    try:
        if file_type == "pdf":
            return _process_pdf(file_path)
        elif file_type in ("png", "jpg", "jpeg", "image"):
            return _process_image(file_path)
        elif file_type == "docx":
            return _process_docx(file_path)
        elif file_type == "pptx":
            return _process_pptx(file_path)
        else:
            logger.warning(f"Unsupported file type: {file_type}")
            return ""
    except Exception as e:
        logger.error(f"Error processing document {file_path}: {e}")
        raise


def _process_pdf(file_path: str) -> str:
    """Extract text from PDF using PyMuPDF, with OCR fallback for image-based pages."""
    doc = fitz.open(file_path)
    text_parts = []

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        text = page.get_text("text")

        if text.strip():
            text_parts.append(f"--- Page {page_num + 1} ---\n{text.strip()}")
        else:
            # Page has no extractable text — skip OCR if not configured
            ocr_path = settings.HUNYUAN_OCR_PATH
            if ocr_path and os.path.exists(ocr_path):
                try:
                    pix = page.get_pixmap(dpi=300)
                    tmp_path = os.path.join(tempfile.gettempdir(), f"ocr_page_{page_num}.png")
                    pix.save(tmp_path)
                    pix = None  # Release pixmap
                    ocr_text = _run_ocr(tmp_path)
                    if ocr_text:
                        text_parts.append(f"--- Page {page_num + 1} (OCR) ---\n{ocr_text}")
                    _safe_remove(tmp_path)
                except Exception as e:
                    logger.warning(f"OCR failed for page {page_num + 1}: {e}")
            else:
                text_parts.append(f"--- Page {page_num + 1} ---\n[Image-based page, OCR not configured]")

    doc.close()
    return "\n\n".join(text_parts)


def _process_image(file_path: str) -> str:
    """Extract text from an image using OCR."""
    return _run_ocr(file_path)


def _process_docx(file_path: str) -> str:
    """Extract text from a Word document."""
    doc = DocxDocument(file_path)
    text_parts = []

    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            text_parts.append(paragraph.text.strip())

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                text_parts.append(row_text)

    return "\n".join(text_parts)


def _process_pptx(file_path: str) -> str:
    """Extract text from a PowerPoint presentation."""
    prs = Presentation(file_path)
    text_parts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        slide_texts.append(paragraph.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        slide_texts.append(row_text)

        if slide_texts:
            text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))

    return "\n\n".join(text_parts)


def _run_ocr(image_path: str) -> str:
    """Run Hunyuan OCR on an image file. Falls back to empty string if OCR is unavailable."""
    ocr_path = settings.HUNYUAN_OCR_PATH
    if not ocr_path or not os.path.exists(ocr_path):
        return ""

    try:
        import subprocess
        result = subprocess.run(
            ["python", os.path.join(ocr_path, "run_ocr.py"), image_path],
            capture_output=True,
            text=True,
            timeout=120,
        )
        if result.returncode == 0:
            return result.stdout.strip()
        else:
            logger.error(f"OCR failed: {result.stderr}")
            return ""
    except Exception as e:
        logger.error(f"OCR error: {e}")
        return ""


def get_page_count(file_path: str, file_type: str) -> int:
    """Get the number of pages in a document."""
    try:
        if file_type == "pdf":
            doc = fitz.open(file_path)
            count = len(doc)
            doc.close()
            return count
        elif file_type == "pptx":
            prs = Presentation(file_path)
            return len(prs.slides)
        elif file_type == "docx":
            return 1
        else:
            return 1
    except Exception:
        return 1


def _safe_remove(path: str):
    """Safely remove a file, ignoring errors on Windows."""
    try:
        if os.path.exists(path):
            os.unlink(path)
    except (PermissionError, OSError):
        pass
